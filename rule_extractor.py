"""
╔══════════════════════════════════════════════════════════════════════════════╗
║  SCHNEIDER ELECTRIC - RULE EXTRACTOR v2.0                                   ║
║  Extract coding rules from PDF, PPT, Word files using AI                   ║
║  Auto-saves to Extracted_Rules_From_Pdf.json + Push to GitHub               ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""

import streamlit as st
import json
import os
import subprocess
from pathlib import Path
from datetime import datetime
import google.generativeai as genai
from dotenv import load_dotenv

# ── File parsers ──────────────────────────────────────────────────────────────
try:
    import pdfplumber
    PDF_OK = True
except ImportError:
    PDF_OK = False

try:
    from pptx import Presentation
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    from docx import Document
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

# ── Config ────────────────────────────────────────────────────────────────────
load_dotenv()

RULES_GUIDE_FOLDER = Path("Rules_Guide_Used")
RULES_OUTPUT_FILE  = Path("server/Extracted_Rules_From_Pdf.json")
GEMINI_API_KEY     = os.getenv("GEMINI_API_KEY", "")
GITHUB_REPO_PATH   = Path(".")

RULES_GUIDE_FOLDER.mkdir(exist_ok=True)

# ── Page setup ────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Rule Extractor | Schneider AI",
    page_icon="📋",
    layout="wide"
)

st.markdown("""
<style>
    .stButton>button { background:#3DCD58; color:white; border:none; border-radius:8px; font-weight:bold; }
    .stButton>button:hover { background:#2db347; }
    .rule-card { background:#1e2130; border-left:4px solid #3DCD58; border-radius:8px; padding:12px 16px; margin:8px 0; }
    .file-card { background:#1a1d2e; border:1px solid #2d3250; border-radius:8px; padding:10px 14px; margin:6px 0; }
    h1 { color:#3DCD58 !important; }
</style>
""", unsafe_allow_html=True)

# ── Gemini setup ──────────────────────────────────────────────────────────────
@st.cache_resource
def get_gemini():
    if not GEMINI_API_KEY:
        return None
    genai.configure(api_key=GEMINI_API_KEY)
    return genai.GenerativeModel("gemini-2.0-flash")

# ── Text extractors ───────────────────────────────────────────────────────────
def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    try:
        if ext == ".pdf":
            if not PDF_OK:
                return "ERROR: run pip install pdfplumber"
            text = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        text.append(t)
            return "\n".join(text)

        elif ext in (".pptx", ".ppt"):
            if not PPTX_OK:
                return "ERROR: run pip install python-pptx"
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text.strip())
            return "\n".join(text)

        elif ext in (".docx", ".doc"):
            if not DOCX_OK:
                return "ERROR: run pip install python-docx"
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif ext == ".txt":
            return open(file_path, encoding="utf-8", errors="ignore").read()

        else:
            return f"ERROR: Unsupported type {ext}"
    except Exception as e:
        return f"ERROR: {e}"

# ── AI extraction ─────────────────────────────────────────────────────────────
def extract_rules_with_ai(text: str, source_file: str, model) -> list:
    prompt = f"""
You are a Schneider Electric coding standards expert.

Read the document below and extract EVERY coding rule, standard, guideline, or best practice.

For each rule return a JSON object with:
- "rule_id": "NEW_001" (placeholder, will be renumbered)
- "rule": clear 1-2 sentence rule statement
- "suggested_fix": how to comply with the rule
- "source": "{source_file}"
- "category": one of [naming, structure, security, energy, documentation, safety, performance, general]
- "severity": one of [critical, error, warning, info]

Return ONLY a valid JSON array. No markdown, no explanation.
If nothing found, return [].

Document:
{text[:8000]}
"""
    try:
        response = model.generate_content(prompt)
        raw = response.text.strip()
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        rules = json.loads(raw.strip())
        return rules if isinstance(rules, list) else []
    except Exception as e:
        st.error(f"AI error: {e}")
        return []

# ── Rules JSON helpers ────────────────────────────────────────────────────────
def load_existing_rules() -> list:
    if RULES_OUTPUT_FILE.exists():
        try:
            return json.loads(RULES_OUTPUT_FILE.read_text(encoding="utf-8")).get("rules", [])
        except Exception:
            return []
    return []

def save_rules(rules: list):
    RULES_OUTPUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    RULES_OUTPUT_FILE.write_text(
        json.dumps({"rules": rules}, indent=2, ensure_ascii=False),
        encoding="utf-8"
    )

def get_next_rule_number(existing: list) -> int:
    nums = ["".join(filter(str.isdigit, r.get("rule_id", ""))) for r in existing]
    nums = [int(n) for n in nums if n]
    return max(nums, default=0) + 1

def renumber(rules: list, start: int) -> list:
    for i, r in enumerate(rules):
        r["rule_id"] = f"R{start + i:03d}"
    return rules

# ── Git helpers ───────────────────────────────────────────────────────────────
def run_git(cmd: list) -> tuple:
    try:
        r = subprocess.run(cmd, cwd=str(GITHUB_REPO_PATH.resolve()),
                           capture_output=True, text=True, timeout=30)
        return r.returncode == 0, r.stdout + r.stderr
    except Exception as e:
        return False, str(e)

def push_to_github(commit_msg: str) -> tuple:
    steps = []

    ok, out = run_git(["git", "add", str(RULES_OUTPUT_FILE)])
    steps.append(("git add", ok, out))
    if not ok:
        return False, steps

    ok, out = run_git(["git", "commit", "-m", commit_msg])
    steps.append(("git commit", ok, out))
    if not ok:
        if "nothing to commit" in out:
            steps.append(("note", True, "Already up to date — no changes to commit."))
            return True, steps
        return False, steps

    ok, out = run_git(["git", "push", "origin", "main"])
    steps.append(("git push", ok, out))
    return ok, steps

def scan_folder() -> list:
    exts = {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".txt"}
    return sorted([f for f in RULES_GUIDE_FOLDER.iterdir()
                   if f.is_file() and f.suffix.lower() in exts])

# ═══════════════════════════════════════════════════════════════════════════════
# UI
# ═══════════════════════════════════════════════════════════════════════════════
st.title("📋 Rule Extractor")
st.markdown("**Upload any PDF / PPT / Word → AI extracts rules → Push to GitHub automatically**")
st.divider()

model = get_gemini()
if not model:
    st.error("⚠️ GEMINI_API_KEY missing from .env")

existing_rules = load_existing_rules()

# Stats bar
c1, c2, c3, c4 = st.columns(4)
c1.metric("📋 Rules in JSON", len(existing_rules))
c2.metric("📁 Files in Folder", len(scan_folder()))
c3.metric("🆕 This Session", len(st.session_state.get("extracted_rules", [])))
c4.metric("🤖 Model", "Gemini 2.0 Flash")
st.divider()

col1, col2 = st.columns(2, gap="large")

# ── LEFT: Upload ──────────────────────────────────────────────────────────────
with col1:
    st.subheader("📁 Upload Files")
    uploaded = st.file_uploader(
        "Drop PDF / PPT / Word / TXT files here",
        type=["pdf", "pptx", "ppt", "docx", "doc", "txt"],
        accept_multiple_files=True,
        label_visibility="collapsed"
    )
    if uploaded:
        for f in uploaded:
            dest = RULES_GUIDE_FOLDER / f.name
            dest.write_bytes(f.read())
            st.success(f"✅ Saved: **{f.name}**")

    st.divider()
    st.subheader("📂 Rules_Guide_Used")
    files = scan_folder()
    if not files:
        st.info("No files yet. Upload above or paste files into the folder.")
    else:
        for fp in files:
            ext = fp.suffix.upper().lstrip(".")
            icon = {"PDF":"📄","PPTX":"📊","PPT":"📊","DOCX":"📝","DOC":"📝","TXT":"📃"}.get(ext,"📄")
            kb = max(1, fp.stat().st_size // 1024)
            st.markdown(
                f'<div class="file-card">{icon} <b>{fp.name}</b> '
                f'<span style="color:#888;float:right">{kb} KB</span></div>',
                unsafe_allow_html=True
            )

# ── RIGHT: Extract ────────────────────────────────────────────────────────────
with col2:
    st.subheader("🤖 Extract with AI")
    files = scan_folder()

    if not files:
        st.info("Upload files on the left first.")
    else:
        file_names = [f.name for f in files]
        selected = st.multiselect("Select files:", options=file_names, default=file_names)

        ca, cb = st.columns(2)
        auto_save  = ca.toggle("Auto-save to JSON", value=True)
        skip_dupes = cb.toggle("Skip duplicates",   value=True)

        if st.button("🚀 Extract Rules Now", use_container_width=True):
            if not model:
                st.error("Need GEMINI_API_KEY")
            elif not selected:
                st.warning("Select at least one file")
            else:
                all_new = []
                prog = st.progress(0, text="Starting...")

                for idx, fname in enumerate(selected):
                    prog.progress(idx / len(selected), text=f"📖 {fname}")
                    with st.spinner(f"Reading {fname}..."):
                        text = extract_text(str(RULES_GUIDE_FOLDER / fname))
                        if text.startswith("ERROR"):
                            st.warning(f"⚠️ {text}")
                            continue
                        rules = extract_rules_with_ai(text, fname, model)
                        if rules:
                            all_new.extend(rules)
                            st.success(f"✅ **{fname}** → {len(rules)} rules")
                        else:
                            st.warning(f"⚠️ No rules found in {fname}")

                prog.progress(1.0, text="✅ Done!")

                if all_new:
                    all_new = renumber(all_new, get_next_rule_number(existing_rules))
                    if skip_dupes:
                        exist_set = {r.get("rule","").lower() for r in existing_rules}
                        all_new = [r for r in all_new if r.get("rule","").lower() not in exist_set]

                    st.session_state["extracted_rules"] = all_new
                    st.info(f"📋 **{len(all_new)} unique new rules** ready")

                    if auto_save:
                        save_rules(existing_rules + all_new)
                        st.session_state["rules_saved"] = True
                        st.success(f"💾 Saved! Total: **{len(existing_rules)+len(all_new)}** rules")
                        st.balloons()
                else:
                    st.error("No rules extracted.")

# ── GITHUB PUSH ───────────────────────────────────────────────────────────────
st.divider()
st.subheader("🐙 Push to GitHub")

extracted = st.session_state.get("extracted_rules", [])

if not extracted:
    st.info("Extract rules first, then push here.")
else:
    default_msg = (f"Add {len(extracted)} new rules from Rules_Guide_Used "
                   f"[{datetime.now().strftime('%Y-%m-%d %H:%M')}]")

    col_msg, col_btn = st.columns([3, 1])
    commit_msg = col_msg.text_input("Commit message:", value=default_msg,
                                    label_visibility="collapsed")

    if col_btn.button("⬆️ Push to GitHub", use_container_width=True):
        if not st.session_state.get("rules_saved"):
            save_rules(existing_rules + extracted)

        with st.spinner("Pushing..."):
            success, steps = push_to_github(commit_msg)

        for name, ok, out in steps:
            icon = "✅" if ok else "❌"
            with st.expander(f"{icon} {name}", expanded=not ok):
                st.code(out or "No output")

        if success:
            st.success("🎉 **Pushed to GitHub successfully!**")
            st.markdown(
                "🔗 [View on GitHub](https://github.com/ShriHarsan64K/Schneider-AI-Code-Reviewer"
                "/blob/main/server/Extracted_Rules_From_Pdf.json)"
            )
        else:
            st.error("❌ Push failed — check details above.")
            st.info("""
**Quick fixes:**
- Run `git config --global user.email "you@email.com"` in terminal
- Make sure you're inside the repo folder when running streamlit
- Check internet connection
""")

# ── PREVIEW ───────────────────────────────────────────────────────────────────
st.divider()
st.subheader("👀 Preview Extracted Rules")

if extracted:
    cats = ["All"] + sorted({r.get("category","general") for r in extracted})
    sel_cat = st.selectbox("Filter by category:", cats)
    filtered = extracted if sel_cat == "All" else [r for r in extracted if r.get("category") == sel_cat]

    st.markdown(f"Showing **{len(filtered)}** rules:")
    for rule in filtered[:50]:
        sev = rule.get("severity","info")
        col = {"critical":"#ff4444","error":"#ff8800","warning":"#ffcc00","info":"#3DCD58"}.get(sev,"#888")
        st.markdown(f"""
<div class="rule-card">
    <div style="display:flex;justify-content:space-between;margin-bottom:4px">
        <b style="color:#3DCD58">{rule.get('rule_id','')}</b>
        <span style="color:{col};font-size:12px;text-transform:uppercase">{sev}</span>
    </div>
    <div style="color:#e0e0e0;margin-bottom:6px">{rule.get('rule','')}</div>
    <div style="color:#888;font-size:13px">🔧 {rule.get('suggested_fix','')}</div>
    <div style="color:#555;font-size:11px;margin-top:4px">📁 {rule.get('source','')} · 🏷️ {rule.get('category','')}</div>
</div>""", unsafe_allow_html=True)
else:
    st.info("Run an extraction above to see rules here.")
